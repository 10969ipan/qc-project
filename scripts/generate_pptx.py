from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Create presentation 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Modern Color Palette
DARK_NAVY = RGBColor(15, 23, 42)        # #0F172A
PRIMARY_BLUE = RGBColor(59, 130, 246)    # #3B82F6
LIGHT_BLUE = RGBColor(96, 165, 250)      # #60A5FA
ACCENT_CYAN = RGBColor(34, 211, 238)     # #22D3EE
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(241, 245, 249)     # #F1F5F9
DARK_GRAY = RGBColor(71, 85, 105)        # #475569
SUCCESS_GREEN = RGBColor(34, 197, 94)    # #22C55E
WARNING_ORANGE = RGBColor(249, 115, 22)  # #F97316

def add_gradient_bg(slide):
    """Add dark gradient background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_accent_bar(slide, top=Inches(0)):
    """Add accent bar at top"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_CYAN
    shape.line.fill.background()

def add_side_accent(slide):
    """Add left side accent"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.fill.background()

def add_corner_decoration(slide):
    """Add corner decoration"""
    # Top right circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(-1), Inches(3), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.fill.fore_color.brightness = 0.3
    shape.line.fill.background()

def add_title_slide(title, subtitle=""):
    """Modern title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_corner_decoration(slide)
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(52)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.paragraphs[0].text = subtitle
        tf2.paragraphs[0].font.size = Pt(24)
        tf2.paragraphs[0].font.color.rgb = LIGHT_BLUE

def add_content_slide(title, bullets, has_screenshot=False):
    """Content slide with optional screenshot area"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_side_accent(slide)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    if has_screenshot:
        # Content on left (narrower)
        content_width = Inches(5.5)
        content_left = Inches(0.8)
        
        # Screenshot placeholder on right
        ss_left = Inches(6.8)
        ss_top = Inches(1.5)
        ss_width = Inches(6)
        ss_height = Inches(5.2)
        
        # Screenshot frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ss_left, ss_top, ss_width, ss_height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(30, 41, 59)  # Darker than bg
        frame.line.color.rgb = PRIMARY_BLUE
        frame.line.width = Pt(2)
        
        # Screenshot label
        label = slide.shapes.add_textbox(ss_left, ss_top + Inches(2.2), ss_width, Inches(0.6))
        label_tf = label.text_frame
        label_tf.paragraphs[0].text = "📷 Screenshot"
        label_tf.paragraphs[0].font.size = Pt(20)
        label_tf.paragraphs[0].font.color.rgb = DARK_GRAY
        label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Hint text
        hint = slide.shapes.add_textbox(ss_left, ss_top + Inches(2.7), ss_width, Inches(0.5))
        hint_tf = hint.text_frame
        hint_tf.paragraphs[0].text = "(Klik kanan → Change Picture)"
        hint_tf.paragraphs[0].font.size = Pt(12)
        hint_tf.paragraphs[0].font.color.rgb = DARK_GRAY
        hint_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        content_width = Inches(11.5)
        content_left = Inches(0.8)
    
    # Content bullets
    txBox2 = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        
        # Check if it's a sub-item
        if bullet.startswith("   "):
            p.text = "    ○ " + bullet.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = LIGHT_BLUE
        else:
            p.text = "● " + bullet
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
        p.space_after = Pt(14)

def add_screenshot_slide(title, description=""):
    """Full screenshot slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_side_accent(slide)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Large screenshot area
    ss_left = Inches(0.8)
    ss_top = Inches(1.4)
    ss_width = Inches(11.7)
    ss_height = Inches(5.5)
    
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ss_left, ss_top, ss_width, ss_height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(30, 41, 59)
    frame.line.color.rgb = PRIMARY_BLUE
    frame.line.width = Pt(2)
    
    # Center label
    label = slide.shapes.add_textbox(ss_left, ss_top + Inches(2.3), ss_width, Inches(0.6))
    label_tf = label.text_frame
    label_tf.paragraphs[0].text = "📷 " + (description if description else "Insert Screenshot Here")
    label_tf.paragraphs[0].font.size = Pt(24)
    label_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    hint = slide.shapes.add_textbox(ss_left, ss_top + Inches(2.9), ss_width, Inches(0.5))
    hint_tf = hint.text_frame
    hint_tf.paragraphs[0].text = "(Klik kanan shape → Change Picture untuk menambahkan gambar)"
    hint_tf.paragraphs[0].font.size = Pt(14)
    hint_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    hint_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_table_slide(title, headers, rows):
    """Modern table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_side_accent(slide)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    left = Inches(0.8)
    top = Inches(1.4)
    width = Inches(11.7)
    row_height = min(0.5, 5.5 / num_rows)
    height = Inches(row_height * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    # Style headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Style data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(51, 65, 85)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_workflow_slide():
    """Visual workflow slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_side_accent(slide)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "ALUR APPROVAL"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Workflow boxes
    box_width = Inches(2.8)
    box_height = Inches(1.4)
    arrow_top = Inches(3.5)
    box_top = Inches(3)
    
    positions = [
        (Inches(0.8), "INSPECTOR\nInput Data", PRIMARY_BLUE),
        (Inches(3.9), "KARU QC\nApproval 1", LIGHT_BLUE),
        (Inches(7), "SUPERVISOR\nApproval 2", ACCENT_CYAN),
        (Inches(10.1), "CLOSE\n✓ Selesai", SUCCESS_GREEN),
    ]
    
    for i, (left, text, color) in enumerate(positions):
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Text in box
        txb = slide.shapes.add_textbox(left, box_top + Inches(0.35), box_width, box_height)
        tff = txb.text_frame
        tff.paragraphs[0].text = text
        tff.paragraphs[0].font.size = Pt(18)
        tff.paragraphs[0].font.bold = True
        tff.paragraphs[0].font.color.rgb = WHITE if color != ACCENT_CYAN else DARK_NAVY
        tff.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow (except last)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + box_width + Inches(0.1), arrow_top, Inches(0.5), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = WHITE
            arrow.line.fill.background()
    
    # Reject flow
    reject_top = Inches(5)
    reject = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), reject_top, Inches(4), Inches(0.8))
    reject.fill.solid()
    reject.fill.fore_color.rgb = WARNING_ORANGE
    reject.line.fill.background()
    
    reject_txt = slide.shapes.add_textbox(Inches(4.5), reject_top + Inches(0.2), Inches(4), Inches(0.6))
    reject_tf = reject_txt.text_frame
    reject_tf.paragraphs[0].text = "❌ REJECT → Revisi Data"
    reject_tf.paragraphs[0].font.size = Pt(16)
    reject_tf.paragraphs[0].font.bold = True
    reject_tf.paragraphs[0].font.color.rgb = WHITE
    reject_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_tips_slide():
    """Tips slide with icons"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_side_accent(slide)
    add_accent_bar(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "TIPS PENGGUNAAN"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Do's column
    do_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.3))
    do_box.fill.solid()
    do_box.fill.fore_color.rgb = RGBColor(20, 83, 45)  # Dark green
    do_box.line.fill.background()
    
    do_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6))
    do_title_tf = do_title.text_frame
    do_title_tf.paragraphs[0].text = "✅ LAKUKAN"
    do_title_tf.paragraphs[0].font.size = Pt(24)
    do_title_tf.paragraphs[0].font.bold = True
    do_title_tf.paragraphs[0].font.color.rgb = SUCCESS_GREEN
    do_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    do_items = ["Login dengan akun sendiri", "Input data segera setelah cek", "Double-check sebelum simpan", "Logout setelah selesai", "Update status mesin/meja"]
    do_content = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(5), Inches(4))
    do_tf = do_content.text_frame
    for i, item in enumerate(do_items):
        if i == 0:
            p = do_tf.paragraphs[0]
        else:
            p = do_tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Don'ts column
    dont_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.4), Inches(5.5), Inches(5.3))
    dont_box.fill.solid()
    dont_box.fill.fore_color.rgb = RGBColor(127, 29, 29)  # Dark red
    dont_box.line.fill.background()
    
    dont_title = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.6))
    dont_title_tf = dont_title.text_frame
    dont_title_tf.paragraphs[0].text = "❌ JANGAN"
    dont_title_tf.paragraphs[0].font.size = Pt(24)
    dont_title_tf.paragraphs[0].font.bold = True
    dont_title_tf.paragraphs[0].font.color.rgb = RGBColor(248, 113, 113)
    dont_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    dont_items = ["Berbagi password", "Tinggalkan browser terbuka", "Input data belum diverifikasi", "Hapus data tanpa konfirmasi"]
    dont_content = slide.shapes.add_textbox(Inches(7.4), Inches(2.3), Inches(5), Inches(4))
    dont_tf = dont_content.text_frame
    for i, item in enumerate(dont_items):
        if i == 0:
            p = dont_tf.paragraphs[0]
        else:
            p = dont_tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

# ==================== CREATE SLIDES ====================

# 1. Title Slide
add_title_slide("SOSIALISASI QC APPS", "Plant Jakarta • 2026")

# 2. Agenda
add_content_slide("AGENDA", [
    "Cara Login ke Aplikasi",
    "Menu untuk Inspector",
    "Menu untuk Karu QC", 
    "Menu untuk Supervisor",
    "Alur Approval",
    "Tips & Troubleshooting"
])

# 3. Login dengan screenshot
add_content_slide("CARA LOGIN", [
    "Buka browser Chrome/Firefox",
    "Akses URL aplikasi",
    "Masukkan Email & Password",
    "Klik tombol LOGIN"
], has_screenshot=True)

# 4. Screenshot halaman login
add_screenshot_slide("TAMPILAN HALAMAN LOGIN", "Screenshot halaman login")

# 5. Daftar Akun
add_table_slide("DAFTAR AKUN PLANT JAKARTA",
    ["Nama", "Email", "Role"],
    [
        ["Masuli", "masuli.jkt@qc.com", "Supervisor"],
        ["Marsiah", "marsiah.jkt@qc.com", "Karu QC"],
        ["Afrin Wibowo", "afrin.jkt@qc.com", "Inspector"],
        ["Anggriyani", "anggriyani.jkt@qc.com", "Inspector"],
        ["Okah Retno A.", "okah.jkt@qc.com", "Inspector"],
        ["M. Miftahul U.", "ulum.jkt@qc.com", "Inspector"],
        ["Ilham Aldi P.", "ilham.jkt@qc.com", "Inspector"],
        ["Tri Rahmadhani", "tri.jkt@qc.com", "Inspector"],
        ["Sabrina K.", "sabrina.jkt@qc.com", "Inspector"],
        ["Ririn Eka P.", "ririn.jkt@qc.com", "Inspector"],
        ["Syadina Juhro", "syadina.jkt@qc.com", "Inspector"],
    ]
)

# 6. Dashboard
add_content_slide("DASHBOARD", [
    "Status Approval checksheet",
    "Produksi Sub Assy (status meja)",
    "Produksi Injection (status mesin)",
    "Monitoring real-time"
], has_screenshot=True)

# 7. Screenshot Dashboard
add_screenshot_slide("TAMPILAN DASHBOARD", "Screenshot dashboard")

# 8. Menu Inspector
add_content_slide("MENU INSPECTOR", [
    "Dashboard → Monitoring produksi",
    "Checksheet → Input data QC:",
    "   Sub Assy",
    "   Inprocess",
    "   Sortir",
    "Laporan → Lihat history data"
], has_screenshot=True)

# 9. Cara Input Checksheet
add_content_slide("CARA INPUT CHECKSHEET", [
    "Pilih menu Checksheet → Plant Jakarta",
    "Pilih jenis: Sub Assy / Inprocess / Sortir",
    "Ketik Kode SAP → Item otomatis terpilih",
    "Isi hasil pengecekan",
    "Klik tombol SIMPAN"
], has_screenshot=True)

# 10. Screenshot Input Sub Assy
add_screenshot_slide("INPUT SUB ASSY", "Screenshot form input Sub Assy")

# 11. Screenshot Input Inprocess
add_screenshot_slide("INPUT INPROCESS", "Screenshot form input Inprocess")

# 12. Menu Karu QC
add_content_slide("MENU KARU QC", [
    "Semua fitur Inspector",
    "PLUS Approval Level 1:",
    "   Buka Laporan → Pilih checksheet",
    "   Klik ✓ untuk Approve",
    "   Klik ✗ untuk Reject"
], has_screenshot=True)

# 13. Screenshot Approval
add_screenshot_slide("TAMPILAN APPROVAL", "Screenshot halaman approval")

# 14. Menu Supervisor
add_content_slide("MENU SUPERVISOR", [
    "Semua fitur Karu QC",
    "PLUS:",
    "   Master Data (Item, Kategori)",
    "   Report Analisis NG",
    "   Claim Customer",
    "   Final Approval (Level 2)"
], has_screenshot=True)

# 15. Alur Approval (Visual)
add_workflow_slide()

# 16. Tips
add_tips_slide()

# 17. Troubleshooting
add_table_slide("TROUBLESHOOTING",
    ["Masalah", "Solusi"],
    [
        ["Tidak bisa login", "Cek email/password, hubungi Admin"],
        ["Halaman error", "Refresh browser (F5)"],
        ["Data tidak tersimpan", "Cek koneksi internet"],
        ["Session expired", "Login ulang"],
    ]
)

# 18. Closing
add_title_slide("TERIMA KASIH", "QC Apps • Plant Jakarta")

# Save
output_path = r"c:\laragon\www\qc-project\sosialisasi_qc_apps_jakarta_v2.pptx"
prs.save(output_path)
print(f"✅ PPTX saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
